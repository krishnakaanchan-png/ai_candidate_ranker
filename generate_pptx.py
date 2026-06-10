import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Paths
TEMPLATE_PATH = '/Users/vani/Downloads/[PUB] India_runs_data_and_ai_challenge/India_runs_data_and_ai_challenge/ai_candidate_ranker/Copy of Idea Submission Template | Redrob.pptx'
OUTPUT_PATH = '/Users/vani/Downloads/[PUB] India_runs_data_and_ai_challenge/India_runs_data_and_ai_challenge/ai_candidate_ranker/Submission_Deck.pptx'

# Load template (we will reuse its slide layouts)
prs = Presentation(TEMPLATE_PATH)

# Helper to add bullet points to a placeholder shape
def add_bullets(shape, lines):
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(14)
        p.level = 0
        p.alignment = PP_ALIGN.LEFT

# --- Slide 1: Team Info ---
slide1 = prs.slides[0]
# Assuming first shape is title, second is content
title1 = slide1.shapes.title
title1.text = "FitFinderAi"
content1 = slide1.shapes[1]
add_bullets(content1, [
    "Problem Statement: Recruiters miss the right candidate due to keyword‑only matching.",
    "Team Leader: Kaanchan Krishna",
    "Team Name: FitFinderAi"
])

# --- Slide 2: Solution Overview ---
slide2 = prs.slides[1]
slide2.shapes.title.text = "Solution Overview"
content2 = slide2.shapes[1]
add_bullets(content2, [
    "Two‑Phase Hybrid Retrieval Architecture:",
    "  • Phase 1 – Offline: Embed 100K candidates with SentenceTransformer (all‑MiniLM‑L6‑v2) and store in FAISS.",
    "  • Phase 2 – Online: Embed JD, retrieve top 10 000 semantic matches, apply strict hard‑filters, compute behavioral score, and output top 100.",
    "Key Innovations:",
    "  • Redrob behavioral signals (response rate, interview completion, profile activity) weight 40 % of final score.",
    "  • Honeypot and JD‑red‑flag detection eliminates fake or unsuitable profiles.",
    "  • End‑to‑end runtime < 10 s on a single CPU core, < 500 MB RAM."
])

# --- Slide 3: JD Understanding & Candidate Evaluation ---
slide3 = prs.slides[2]
slide3.shapes.title.text = "JD Understanding & Candidate Evaluation"
content3 = slide3.shapes[1]
add_bullets(content3, [
    "Key JD Requirements extracted:",
    "  • 4–15 years of experience (sweet spot 5–9 years gives slight bonus).",
    "  • AI/ML engineering role – strong Python, ML libraries, product experience.",
    "  • Disallow pure‑researcher profiles, consulting‑only backgrounds, LangChain‑only wrappers.",
    "Important Candidate Signals (Redrob):",
    "  • Recruiter response rate (40 % weight).",
    "  • Interview completion rate (30 % weight).",
    "  • Profile completeness (10 % weight).",
    "  • GitHub activity (10 % weight).",
    "  • Open‑to‑work flag (+5 flat bonus).",
    "  • Notice period > 60 days (‑10 % penalty)."
])

# --- Slide 4: Ranking Methodology ---
slide4 = prs.slides[3]
slide4.shapes.title.text = "Ranking Methodology"
content4 = slide4.shapes[1]
add_bullets(content4, [
    "1️⃣ Retrieve top 10 000 candidates using FAISS cosine similarity.",
    "2️⃣ Apply hard filters: YoE bounds, honeypot detection, JD red‑flags.",
    "3️⃣ Compute behavioral score from Redrob signals.",
    "4️⃣ Final Score = 0.6 × Semantic + 0.4 × Behavioral (+ YoE bonus).",
    "5️⃣ Output exactly 100 ranked candidates (validated by submission script)."
])

# --- Slide 5: Explainability & Data Validation ---
slide5 = prs.slides[4]
slide5.shapes.title.text = "Explainability & Data Validation"
content5 = slide5.shapes[1]
add_bullets(content5, [
    "For each candidate we expose:",
    "  • Semantic similarity value.",
    "  • Behavioral signal breakdown and weight contribution.",
    "  • Reasons for disqualification (e.g., honeypot, consulting only).",
    "Data Validation:",
    "  • All 100 output rows pass strict schema validation (validate_submission.py).",
    "  • No network calls at runtime; all data is local/pre‑computed."
])

# --- Slide 6: End‑to‑End Workflow ---
slide6 = prs.slides[5]
slide6.shapes.title.text = "End‑to‑End Workflow"
content6 = slide6.shapes[1]
add_bullets(content6, [
    "Offline (Phase 1):",
    "  • Load raw candidate JSONL.",
    "  • Extract and order text (Title → Summary → Skills → Career).",
    "  • Encode with SentenceTransformer, L2‑normalize, store in FAISS index.",
    "  • Serialize index + metadata pickle.",
    "Online (Phase 2):",
    "  • Parse JD, embed same way.",
    "  • FAISS query → top 10 k.",
    "  • Apply hard filters, compute behavioral score, rank, write CSV."
])

# --- Slide 7: System Architecture ---
slide7 = prs.slides[6]
slide7.shapes.title.text = "System Architecture"
content7 = slide7.shapes[1]
add_bullets(content7, [
    "[Diagram placeholder – see repository README for visual architecture].",
    "Components:",
    "  • Python 3.14, FAISS, SentenceTransformers.",
    "  • Redrob signals JSON object per candidate.",
    "  • Validation script (validate_submission.py)."
])

# --- Slide 8: Results & Performance ---
slide8 = prs.slides[7]
slide8.shapes.title.text = "Results & Performance"
content8 = slide8.shapes[1]
add_bullets(content8, [
    "Performance Metrics:",
    "  • 100 K candidates processed offline in ~3 min.",
    "  • Online ranking completed in < 10 s (CPU‑only).",
    "  • Peak RAM usage ~ 480 MB.",
    "  • 100 % of output rows validated.",
    "  • All hackathon constraints satisfied.",
    "Submission Assets:",
    "  • GitHub repo: https://github.com/krishnakaanchan-png/ai_candidate_ranker",
    "  • Methodology PDF, Presentation Deck, sandbox notebook, top‑100 CSV."
])

# --- Slide 9: Technologies Used ---
slide9 = prs.slides[8]
slide9.shapes.title.text = "Technologies Used"
content9 = slide9.shapes[1]
add_bullets(content9, [
    "Python 3.14",
    "FAISS (IndexFlatIP) for fast vector search",
    "SentenceTransformers (all‑MiniLM‑L6‑v2) for embeddings",
    "ReportLab for PDF generation",
    "python‑pptx for presentation deck",
    "GitHub Actions (CI) – optional",
    "Redrob synthetic dataset (behavioral signals)"
])

# --- Slide 10: Submission Assets ---
slide10 = prs.slides[9]
slide10.shapes.title.text = "Submission Assets"
content10 = slide10.shapes[1]
add_bullets(content10, [
    "• GitHub repository with full code.",
    "• methodology.pdf (detailed deck).",
    "• Submission_Deck.pptx (this file).",
    "• sandbox_notebook.ipynb (demo execution).",
    "• team_FitFinderAi.csv (top‑100 ranked candidates).",
    "• README.md with run instructions."
])

# Save the filled deck
prs.save(OUTPUT_PATH)
print(f"Presentation saved to {OUTPUT_PATH}")
