import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

TEMPLATE_PATH = '/Users/vani/Downloads/[PUB] India_runs_data_and_ai_challenge/India_runs_data_and_ai_challenge/ai_candidate_ranker/Copy of Idea Submission Template | Redrob.pptx'
OUTPUT_PATH = '/Users/vani/Downloads/[PUB] India_runs_data_and_ai_challenge/India_runs_data_and_ai_challenge/ai_candidate_ranker/Submission_Deck_Filled.pptx'

prs = Presentation(TEMPLATE_PATH)

def set_text(shape, lines, font_size=14):
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(font_size)
        p.alignment = PP_ALIGN.LEFT

def get_title_and_content(slide):
    # title is usually shape index 1 (after picture)
    title = slide.shapes[1] if len(slide.shapes) > 1 else None
    # content placeholder may be index 2; if not present, create a textbox
    if len(slide.shapes) > 2:
        content = slide.shapes[2]
    else:
        # add a textbox covering most of the slide (leaving some margin)
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        content = slide.shapes.add_textbox(left, top, width, height)
    return title, content

# Slide 1 – Team Info
title1, content1 = get_title_and_content(prs.slides[0])
if title1: title1.text = "FitFinderAi"
set_text(content1, [
    "Problem Statement: Recruiters miss the right candidate due to keyword‑only matching.",
    "Team Leader: Kaanchan Krishna",
    "Team Name: FitFinderAi"
])

# Slide 2 – Solution Overview
title2, content2 = get_title_and_content(prs.slides[1])
if title2: title2.text = "Solution Overview"
set_text(content2, [
    "Two‑Phase Hybrid Retrieval Architecture:",
    "  • Phase 1 – Offline: Embed 100K candidates with SentenceTransformer (all‑MiniLM‑L6‑v2) and store in FAISS.",
    "  • Phase 2 – Online: Embed JD, retrieve top 10 000 semantic matches, apply strict hard‑filters, compute behavioral score, and output top 100.",
    "Key Innovations:",
    "  • Redrob behavioral signals (response rate, interview completion, profile activity) weight 40 % of final score.",
    "  • Honeypot and JD‑red‑flag detection eliminates fake or unsuitable profiles.",
    "  • End‑to‑end runtime < 10 s on a single CPU core, < 500 MB RAM."
])

# Slide 3 – JD Understanding & Candidate Evaluation
title3, content3 = get_title_and_content(prs.slides[2])
if title3: title3.text = "JD Understanding & Candidate Evaluation"
set_text(content3, [
    "Key JD Requirements extracted:",
    "  • 4–15 years of experience (sweet spot 5–9 years gives slight bonus).",
    "  • AI/ML engineering role – strong Python, ML libraries, product experience.",
    "  • Disallow pure‑researcher profiles, consulting‑only backgrounds, LangChain‑only wrappers.",
    "Important Candidate Signals (Redrob):",
    "  • Recruiter response rate (40 % weight).",
    "  • Interview completion rate (30 % weight).",
    "  • Profile completeness (10 % weight).",
    "  • GitHub activity (10 % weight).",
    "  • Open‑to‑work flag (+5 flat bonus).",
    "  • Notice period > 60 days (‑10 % penalty)."
])

# Slide 4 – Ranking Methodology
title4, content4 = get_title_and_content(prs.slides[3])
if title4: title4.text = "Ranking Methodology"
set_text(content4, [
    "1️⃣ Retrieve top 10 000 candidates using FAISS cosine similarity.",
    "2️⃣ Apply hard filters: YoE bounds, honeypot detection, JD red‑flags.",
    "3️⃣ Compute behavioral score from Redrob signals.",
    "4️⃣ Final Score = 0.6 × Semantic + 0.4 × Behavioral (+ YoE bonus).",
    "5️⃣ Output exactly 100 ranked candidates (validated by submission script)."
])

# Slide 5 – Explainability & Data Validation
title5, content5 = get_title_and_content(prs.slides[4])
if title5: title5.text = "Explainability & Data Validation"
set_text(content5, [
    "For each candidate we expose:",
    "  • Semantic similarity value.",
    "  • Behavioral signal breakdown and weight contribution.",
    "  • Reasons for disqualification (e.g., honeypot, consulting only).",
    "Data Validation:",
    "  • All 100 output rows pass strict schema validation (validate_submission.py).",
    "  • No network calls at runtime; all data is local/pre‑computed."
])

# Slide 6 – End‑to‑End Workflow
title6, content6 = get_title_and_content(prs.slides[5])
if title6: title6.text = "End‑to‑End Workflow"
set_text(content6, [
    "Offline (Phase 1):",
    "  • Load raw candidate JSONL.",
    "  • Extract and order text (Title → Summary → Skills → Career).",
    "  • Encode with SentenceTransformer, L2‑normalize, store in FAISS index.",
    "  • Serialize index + metadata pickle.",
    "Online (Phase 2):",
    "  • Parse JD, embed same way.",
    "  • FAISS query → top 10 k.",
    "  • Apply hard filters, compute behavioral score, rank, write CSV."
])

# Slide 7 – System Architecture (only title shape exists, add content box)
title7, content7 = get_title_and_content(prs.slides[6])
if title7: title7.text = "System Architecture"
set_text(content7, [
    "Components:",
    "  • Python 3.14, FAISS, SentenceTransformers.",
    "  • Redrob signals JSON object per candidate.",
    "  • Validation script (validate_submission.py).",
    "  • ReportLab for PDF generation, python‑pptx for PPT.",
    "Diagram (see README for visual): offline embedding pipeline → online ranking service."
])

# Slide 8 – Results & Performance
title8, content8 = get_title_and_content(prs.slides[7])
if title8: title8.text = "Results & Performance"
set_text(content8, [
    "Performance Metrics:",
    "  • 100 K candidates processed offline in ~3 min.",
    "  • Online ranking completed in < 10 s (CPU‑only).",
    "  • Peak RAM usage ~ 480 MB.",
    "  • 100 % of output rows validated.",
    "  • All hackathon constraints satisfied.",
    "Submission Assets:",
    "  • GitHub repo: https://github.com/krishnakaanchan-png/ai_candidate_ranker",
    "  • methodology.pdf, Submission_Deck_Filled.pptx, sandbox_notebook.ipynb, top‑100 CSV."
])

# Save filled presentation
prs.save(OUTPUT_PATH)
print(f"Presentation saved to {OUTPUT_PATH}")
